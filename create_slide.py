from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
import re

def clean_text(text):
    if pd.isna(text):
        return ""
    text = str(text).strip()
    return text.replace('\u200b', '')

def create_demo_slide():
    input_file = 'Demo Setup.xlsx'
    output_pptx = 'Cives_Demo_Overview.pptx'

    # Read data to find categories
    try:
        df = pd.read_excel(input_file, header=None)
        # Extract Category column (Col 3 / D)
        raw_cats = df.iloc[:, 3].unique()
        categories = []
        for cat in raw_cats:
            c = clean_text(cat)
            if c and c[0].isdigit(): # Simple heuristic: starts with a number "1. Intro"
                categories.append(c)
            elif c and "Introduction" in c: # Fallback if numbering is weird
                 categories.append(c)
            # Add other specific keywords if needed, but the list seemed numbered
    except Exception as e:
        print(f"Error reading excel for slide: {e}")
        categories = ["Introduction", "Pre-Hearing", "Hearings", "Post-Hearing", "Reporting"]

    prs = Presentation()
    
    # Use a blank layout (6)
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "Cives HMCTS Demo Overview"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

    # Add Content (Two columns if many categories)
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(4.5)
    height = Inches(5)

    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf_content = content_box.text_frame
    tf_content.word_wrap = True

    for cat in categories:
        p = tf_content.add_paragraph()
        p.text = f"• {cat}"
        p.font.size = Pt(20)
        p.space_after = Pt(10)

    # Add a "visual" placeholder or footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf_foot = footer_box.text_frame
    tf_foot.text = "Comprehensive Court Management Solution"
    tf_foot.paragraphs[0].font.italic = True
    tf_foot.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

    prs.save(output_pptx)
    print(f"Successfully created {output_pptx}")

if __name__ == "__main__":
    create_demo_slide()
